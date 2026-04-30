import io
import os
import tempfile
import zipfile

import numpy as np
import streamlit as st

from core.impedance import compute_tdr_diff, compute_tdr_single
from core.mixed_mode import (
    get_SCD21,
    get_SDC21,
    get_SDD11,
    get_SDD21,
    get_SDD22,
    single_to_mixed_mode,
)
from core.parser import get_frequency_ghz, get_smatrix, load_s4p
from core.plots import (
    plot_impedance,
    plot_insertion_loss,
    plot_mode_conversion,
    plot_psfext,
    plot_psnext,
    plot_return_loss,
)


_IMG_EXPORT = dict(format="png", width=630, height=450, scale=2)


def _to_db(s: np.ndarray) -> np.ndarray:
    return 20 * np.log10(np.abs(s) + 1e-15)


def _process_file(file_bytes: bytes, mapping: str, rise_time_ps: float, label: str,
                  tmp_path: str) -> tuple:
    net     = load_s4p(tmp_path)
    freq    = get_frequency_ghz(net)
    fhz     = net.f
    s_se    = get_smatrix(net)
    s_mm    = single_to_mixed_mode(s_se, mapping=mapping)

    t_f,  z11_f  = compute_tdr_single(s_se, fhz, rise_time_ps=rise_time_ps, forward=True)
    _,    zdif_f = compute_tdr_diff(s_mm,   fhz, rise_time_ps=rise_time_ps, forward=True)
    t_r,  z11_r  = compute_tdr_single(s_se, fhz, rise_time_ps=rise_time_ps, forward=False)
    _,    zdif_r = compute_tdr_diff(s_mm,   fhz, rise_time_ps=rise_time_ps, forward=False)

    fd = {
        "freq":  freq,
        "sdd21": get_SDD21(s_mm), "sdd11": get_SDD11(s_mm), "sdd22": get_SDD22(s_mm),
        "scd21": get_SCD21(s_mm), "sdc21": get_SDC21(s_mm),
        "label": label,
    }
    td = {
        "t_fwd": t_f,  "z11_fwd":   z11_f,  "zdiff_fwd": zdif_f,
        "t_rev": t_r,  "z11_rev":   z11_r,  "zdiff_rev": zdif_r,
        "label": label,
    }
    return fd, td


@st.cache_data(show_spinner=False)
def _build_csv_zip(files: tuple, mapping: str, rise_time_ps: float) -> bytes:
    import tempfile as _tmp, os as _os
    tmp_paths = []
    try:
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for file_bytes, filename in files:
                with _tmp.NamedTemporaryFile(delete=False, suffix=".s4p") as f:
                    f.write(file_bytes)
                    tmp_paths.append(f.name)

                fd, td = _process_file(file_bytes, mapping, rise_time_ps, filename, tmp_paths[-1])
                stem   = filename.rsplit(".", 1)[0]

                sdd21_db = _to_db(fd["sdd21"]); sdd11_db = _to_db(fd["sdd11"])
                sdd22_db = _to_db(fd["sdd22"]); scd21_db = _to_db(fd["scd21"])
                sdc21_db = _to_db(fd["sdc21"])

                fd_rows = ["Frequency_GHz,SDD21_dB,SDD11_dB,SDD22_dB,SCD21_dB,SDC21_dB"]
                for i, fv in enumerate(fd["freq"]):
                    fd_rows.append(f"{fv:.6f},{sdd21_db[i]:.6f},{sdd11_db[i]:.6f},"
                                   f"{sdd22_db[i]:.6f},{scd21_db[i]:.6f},{sdc21_db[i]:.6f}")

                td_rows = ["Time_ps,Z_SE_Fwd_Ohm,Z_Diff_Fwd_Ohm,Z_SE_Rev_Ohm,Z_Diff_Rev_Ohm"]
                for i, tv in enumerate(td["t_fwd"]):
                    td_rows.append(f"{tv:.6f},{td['z11_fwd'][i]:.6f},{td['zdiff_fwd'][i]:.6f},"
                                   f"{td['z11_rev'][i]:.6f},{td['zdiff_rev'][i]:.6f}")

                zf.writestr(f"{stem}_frequency_domain.csv", "\n".join(fd_rows))
                zf.writestr(f"{stem}_time_domain.csv",      "\n".join(td_rows))
        return buf.getvalue()
    finally:
        for p in tmp_paths:
            try: _os.unlink(p)
            except: pass


@st.cache_data(show_spinner=False)
def _build_export_zip(
    files: tuple,
    mapping: str, rise_time_ps: float,
    show_zdiff: bool, show_zse: bool, show_forward: bool, show_reverse: bool,
    z_xmin: float, z_xmax: float, z_xstep: float,
    z_ymin: float, z_ymax: float, z_ystep: float,
    fd_xmin: float, fd_xmax: float, fd_xstep: float,
    il_ymin: float, il_ymax: float, il_ystep: float,
    rl_ymin: float, rl_ymax: float, rl_ystep: float,
    psnext_ymin: float, psnext_ymax: float, psnext_ystep: float,
    psfext_ymin: float, psfext_ymax: float, psfext_ystep: float,
    mc_ymin: float, mc_ymax: float, mc_ystep: float,
    show_sdd11: bool = True, show_sdd22: bool = True,
    next_files: tuple = (), fext_files: tuple = (),
    ver: int = 20,
) -> bytes:
    import tempfile as _tmp, os as _os
    tmp_paths = []
    try:
        fd_datasets, td_datasets = [], []
        for file_bytes, filename in files:
            with _tmp.NamedTemporaryFile(delete=False, suffix=".s4p") as f:
                f.write(file_bytes)
                tmp_paths.append(f.name)
            fd, td = _process_file(file_bytes, mapping, rise_time_ps, filename, tmp_paths[-1])
            fd_datasets.append(fd)
            td_datasets.append(td)

        next_datasets = []
        for file_bytes, filename in next_files:
            with _tmp.NamedTemporaryFile(delete=False, suffix=".s4p") as f:
                f.write(file_bytes)
                tmp_paths.append(f.name)
            fd, _ = _process_file(file_bytes, mapping, rise_time_ps, filename, tmp_paths[-1])
            next_datasets.append(fd)

        fext_datasets = []
        for file_bytes, filename in fext_files:
            with _tmp.NamedTemporaryFile(delete=False, suffix=".s4p") as f:
                f.write(file_bytes)
                tmp_paths.append(f.name)
            fd, _ = _process_file(file_bytes, mapping, rise_time_ps, filename, tmp_paths[-1])
            fext_datasets.append(fd)

        figs = {
            "01_insertion_loss":  plot_insertion_loss(
                fd_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
                y_min=il_ymin, y_max=il_ymax, y_step=il_ystep,
            ),
            "02_return_loss":     plot_return_loss(
                fd_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
                y_min=rl_ymin, y_max=rl_ymax, y_step=rl_ystep,
                show_sdd11=show_sdd11, show_sdd22=show_sdd22,
            ),
            "03_psnext":          plot_psnext(
                next_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
                y_min=psnext_ymin, y_max=psnext_ymax, y_step=psnext_ystep,
            ),
            "04_psfext":          plot_psfext(
                fext_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
                y_min=psfext_ymin, y_max=psfext_ymax, y_step=psfext_ystep,
            ),
            "05_mode_conversion": plot_mode_conversion(
                fd_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
                y_min=mc_ymin, y_max=mc_ymax, y_step=mc_ystep,
            ),
            "06_tdr_impedance":   plot_impedance(
                td_datasets, rise_time_ps=rise_time_ps,
                show_zdiff=show_zdiff, show_zse=show_zse,
                show_forward=show_forward, show_reverse=show_reverse,
                x_min=z_xmin, x_max=z_xmax, x_step=z_xstep,
                y_min=z_ymin, y_max=z_ymax, y_step=z_ystep,
            ),
        }
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for name, fig in figs.items():
                zf.writestr(f"{name}.png", fig.to_image(**_IMG_EXPORT))
        return buf.getvalue()
    finally:
        for p in tmp_paths:
            try: _os.unlink(p)
            except: pass


@st.cache_data(show_spinner=False)
def _build_pptx(
    files: tuple,
    mapping: str, rise_time_ps: float,
    show_zdiff: bool, show_zse: bool, show_forward: bool, show_reverse: bool,
    z_xmin: float, z_xmax: float, z_xstep: float,
    z_ymin: float, z_ymax: float, z_ystep: float,
    fd_xmin: float, fd_xmax: float, fd_xstep: float,
    il_ymin: float, il_ymax: float, il_ystep: float,
    rl_ymin: float, rl_ymax: float, rl_ystep: float,
    psnext_ymin: float, psnext_ymax: float, psnext_ystep: float,
    psfext_ymin: float, psfext_ymax: float, psfext_ystep: float,
    mc_ymin: float, mc_ymax: float, mc_ystep: float,
    show_sdd11: bool = True, show_sdd22: bool = True,
    next_files: tuple = (), fext_files: tuple = (),
    ver: int = 20,
) -> bytes:
    import tempfile as _tmp, os as _os
    from pptx import Presentation
    tmp_paths = []
    try:
        fd_datasets, td_datasets = [], []
        for file_bytes, filename in files:
            with _tmp.NamedTemporaryFile(delete=False, suffix=".s4p") as f:
                f.write(file_bytes)
                tmp_paths.append(f.name)
            fd, td = _process_file(file_bytes, mapping, rise_time_ps, filename, tmp_paths[-1])
            fd_datasets.append(fd)
            td_datasets.append(td)

        next_datasets = []
        for file_bytes, filename in next_files:
            with _tmp.NamedTemporaryFile(delete=False, suffix=".s4p") as f:
                f.write(file_bytes)
                tmp_paths.append(f.name)
            fd, _ = _process_file(file_bytes, mapping, rise_time_ps, filename, tmp_paths[-1])
            next_datasets.append(fd)

        fext_datasets = []
        for file_bytes, filename in fext_files:
            with _tmp.NamedTemporaryFile(delete=False, suffix=".s4p") as f:
                f.write(file_bytes)
                tmp_paths.append(f.name)
            fd, _ = _process_file(file_bytes, mapping, rise_time_ps, filename, tmp_paths[-1])
            fext_datasets.append(fd)

        # Slide 1: idx 13=IL, 14=RL, 15=PSNEXT, 16=PSFEXT
        # Slide 2: idx 13=Mode Conversion
        # Slide 3: idx 13=TDR Impedance
        slide_figs = [
            {
                13: plot_insertion_loss(fd_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep, y_min=il_ymin, y_max=il_ymax, y_step=il_ystep),
                14: plot_return_loss(fd_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep, y_min=rl_ymin, y_max=rl_ymax, y_step=rl_ystep, show_sdd11=show_sdd11, show_sdd22=show_sdd22),
                15: plot_psnext(next_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep, y_min=psnext_ymin, y_max=psnext_ymax, y_step=psnext_ystep),
                16: plot_psfext(fext_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep, y_min=psfext_ymin, y_max=psfext_ymax, y_step=psfext_ystep),
            },
            {13: plot_mode_conversion(fd_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep, y_min=mc_ymin, y_max=mc_ymax, y_step=mc_ystep)},
            {13: plot_impedance(td_datasets, rise_time_ps=rise_time_ps, show_zdiff=show_zdiff, show_zse=show_zse, show_forward=show_forward, show_reverse=show_reverse, x_min=z_xmin, x_max=z_xmax, x_step=z_xstep, y_min=z_ymin, y_max=z_ymax, y_step=z_ystep)},
        ]

        template_path = os.path.join(os.path.dirname(__file__), "Template.pptx")
        prs = Presentation(template_path)

        for slide, figs in zip(prs.slides, slide_figs):
            for ph in slide.placeholders:
                idx = ph.placeholder_format.idx
                if idx in figs:
                    ph.insert_picture(io.BytesIO(figs[idx].to_image(**_IMG_EXPORT)))

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    finally:
        for p in tmp_paths:
            try: _os.unlink(p)
            except: pass


st.set_page_config(page_title="SI Tool", page_icon="📡", layout="wide")
st.title("SI Tool — S4P Analyzer")

# ── Sidebar ──────────────────────────────────────────────
with st.sidebar:
    st.header("上傳檔案")
    uploaded_files = st.file_uploader("選擇 THRU.S4P 檔案", type=["s4p"], accept_multiple_files=True)
    uploaded_next_files = st.file_uploader("選擇 NEXT.S4P 檔案", type=["s4p"], accept_multiple_files=True)
    uploaded_fext_files = st.file_uploader("選擇 FEXT.S4P 檔案", type=["s4p"], accept_multiple_files=True)

    st.divider()
    st.subheader("Port Mapping")
    mapping_choice = st.radio(
        "Port Mapping",
        ["Odd-Even", "N+1"],
        index=0,
        horizontal=True,
        label_visibility="hidden",
    )
    mapping = 'A' if mapping_choice == 'Odd-Even' else 'B'

    _TABLE_STYLE = """
<style>
.pm-table { width:100%; border-collapse:collapse; text-align:center; font-size:0.85rem; }
.pm-table td { padding: 4px 6px; border: 1px solid #444; }
.pm-diff { font-weight:bold; background:#dbeafe; color:#1e40af; }
</style>
"""
    if mapping == 'A':
        st.markdown(_TABLE_STYLE + """
<table class="pm-table">
  <tr>
    <td class="pm-diff" rowspan="2">Diff 1</td>
    <td>Port 1</td><td>────</td><td>Port 2</td>
    <td class="pm-diff" rowspan="2">Diff 2</td>
  </tr>
  <tr>
    <td>Port 3</td><td>────</td><td>Port 4</td>
  </tr>
</table>
""", unsafe_allow_html=True)
    else:
        st.markdown(_TABLE_STYLE + """
<table class="pm-table">
  <tr>
    <td class="pm-diff" rowspan="2">Diff 1</td>
    <td>Port 1</td><td>────</td><td>Port 3</td>
    <td class="pm-diff" rowspan="2">Diff 2</td>
  </tr>
  <tr>
    <td>Port 2</td><td>────</td><td>Port 4</td>
  </tr>
</table>
""", unsafe_allow_html=True)

    st.divider()
    st.subheader("Frequency Domain")
    cx1, cx2, cx3 = st.columns(3)
    with cx1:
        fd_xmin  = st.number_input("X min (GHz)", value=0.0,  step=1.0)
    with cx2:
        fd_xmax  = st.number_input("X max (GHz)", value=20.0, step=1.0)
    with cx3:
        fd_xstep = st.number_input("X step (GHz)", value=5.0, step=1.0, min_value=0.1)

    st.caption("Insertion Loss")
    iy1, iy2, iy3 = st.columns(3)
    with iy1:
        il_ymin  = st.number_input("Y min", key="il_ymin",  value=-5.0, step=1.0)
    with iy2:
        il_ymax  = st.number_input("Y max", key="il_ymax",  value=0.0,  step=1.0)
    with iy3:
        il_ystep = st.number_input("Y step", key="il_ystep", value=1.0, step=0.5, min_value=0.1)

    st.caption("Return Loss")
    ry1, ry2, ry3 = st.columns(3)
    with ry1:
        rl_ymin  = st.number_input("Y min", key="rl_ymin",  value=-80.0, step=10.0)
    with ry2:
        rl_ymax  = st.number_input("Y max", key="rl_ymax",  value=0.0,   step=10.0)
    with ry3:
        rl_ystep = st.number_input("Y step", key="rl_ystep", value=20.0, step=5.0, min_value=0.1)
    rl_col1, rl_col2 = st.columns(2)
    with rl_col1:
        show_sdd11 = st.checkbox("SDD11", value=True)
    with rl_col2:
        show_sdd22 = st.checkbox("SDD22", value=True)

    st.caption("PSNEXT")
    ny1, ny2, ny3 = st.columns(3)
    with ny1:
        psnext_ymin  = st.number_input("Y min", key="psnext_ymin",  value=-80.0, step=10.0)
    with ny2:
        psnext_ymax  = st.number_input("Y max", key="psnext_ymax",  value=0.0,   step=10.0)
    with ny3:
        psnext_ystep = st.number_input("Y step", key="psnext_ystep", value=20.0, step=5.0, min_value=0.1)

    st.caption("PSFEXT")
    fy1, fy2, fy3 = st.columns(3)
    with fy1:
        psfext_ymin  = st.number_input("Y min", key="psfext_ymin",  value=-80.0, step=10.0)
    with fy2:
        psfext_ymax  = st.number_input("Y max", key="psfext_ymax",  value=0.0,   step=10.0)
    with fy3:
        psfext_ystep = st.number_input("Y step", key="psfext_ystep", value=20.0, step=5.0, min_value=0.1)

    st.caption("Mode Conversion")
    my1, my2, my3 = st.columns(3)
    with my1:
        mc_ymin  = st.number_input("Y min", key="mc_ymin",  value=-80.0, step=10.0)
    with my2:
        mc_ymax  = st.number_input("Y max", key="mc_ymax",  value=0.0,   step=10.0)
    with my3:
        mc_ystep = st.number_input("Y step", key="mc_ystep", value=20.0, step=5.0, min_value=0.1)

    st.divider()
    st.subheader("Time Domain")
    rise_time_ps = st.number_input(
        "Rise time 20%~80% (ps)",
        min_value=10,
        max_value=500,
        value=35,
        step=1,
    )
    st.caption(f"建議頻寬：{0.35 / (rise_time_ps * 1e-12) / 1e9 * 2:.1f} GHz")
    col_xmin, col_xmax, col_xstep = st.columns(3)
    with col_xmin:
        z_xmin  = st.number_input("X min (ns)",  value=0.0, step=0.1, format="%.2f")
    with col_xmax:
        z_xmax  = st.number_input("X max (ns)",  value=1.0, step=0.1, format="%.2f")
    with col_xstep:
        z_xstep = st.number_input("X step (ns)", value=0.2, step=0.1, min_value=0.001, format="%.3f")
    col_l, col_r = st.columns(2)
    with col_l:
        show_zdiff = st.checkbox("Z_Diff", value=True)
    with col_r:
        show_zse = st.checkbox("Z_SE", value=False)
    col_fwd, col_rev = st.columns(2)
    with col_fwd:
        show_forward = st.checkbox("Forward", value=True)
    with col_rev:
        show_reverse = st.checkbox("Reverse", value=False)
    col_min, col_max, col_ystep = st.columns(3)
    with col_min:
        z_ymin = st.number_input("Y min (Ω)", value=0,   step=10)
    with col_max:
        z_ymax = st.number_input("Y max (Ω)", value=150, step=10)
    with col_ystep:
        z_ystep = st.number_input("Y step (Ω)", value=50, step=10, min_value=1)

# ── Main ─────────────────────────────────────────────────
if not uploaded_files:
    st.info("← 請從左側上傳 S4P 檔案")
    st.stop()

fd_datasets: list = []
td_datasets: list = []
tmp_paths:   list = []

try:
    for uploaded in uploaded_files:
        file_bytes = uploaded.getvalue()
        with tempfile.NamedTemporaryFile(delete=False, suffix=".s4p") as tmp:
            tmp.write(file_bytes)
            tmp_paths.append(tmp.name)

        fd, td = _process_file(file_bytes, mapping, rise_time_ps, uploaded.name, tmp_paths[-1])
        fd_datasets.append(fd)
        td_datasets.append(td)

        st.success(
            f"**{uploaded.name}** — {len(fd['freq'])} 點 | "
            f"{fd['freq'][0]:.3f} ~ {fd['freq'][-1]:.3f} GHz"
        )

    with st.sidebar:
        st.divider()
        st.subheader("輸出")
        stem = uploaded_files[0].name.rsplit(".", 1)[0] if len(uploaded_files) == 1 else "combined"
        files_tuple      = tuple((uf.getvalue(), uf.name) for uf in uploaded_files)
        next_files_tuple = tuple((uf.getvalue(), uf.name) for uf in (uploaded_next_files or []))
        fext_files_tuple = tuple((uf.getvalue(), uf.name) for uf in (uploaded_fext_files or []))

        try:
            with st.spinner("生成圖片中…"):
                zip_bytes = _build_export_zip(
                    files_tuple, mapping, rise_time_ps,
                    show_zdiff, show_zse, show_forward, show_reverse,
                    z_xmin, z_xmax, z_xstep, z_ymin, z_ymax, z_ystep,
                    fd_xmin, fd_xmax, fd_xstep,
                    il_ymin, il_ymax, il_ystep,
                    rl_ymin, rl_ymax, rl_ystep,
                    psnext_ymin, psnext_ymax, psnext_ystep,
                    psfext_ymin, psfext_ymax, psfext_ystep,
                    mc_ymin, mc_ymax, mc_ystep,
                    show_sdd11=show_sdd11, show_sdd22=show_sdd22,
                    next_files=next_files_tuple, fext_files=fext_files_tuple,
                )
            st.download_button(
                "📥 下載全部圖表 (PNG)",
                zip_bytes,
                file_name=f"{stem}_charts.zip",
                mime="application/zip",
            )

            with st.spinner("生成 PPTX 中…"):
                pptx_bytes = _build_pptx(
                    files_tuple, mapping, rise_time_ps,
                    show_zdiff, show_zse, show_forward, show_reverse,
                    z_xmin, z_xmax, z_xstep, z_ymin, z_ymax, z_ystep,
                    fd_xmin, fd_xmax, fd_xstep,
                    il_ymin, il_ymax, il_ystep,
                    rl_ymin, rl_ymax, rl_ystep,
                    psnext_ymin, psnext_ymax, psnext_ystep,
                    psfext_ymin, psfext_ymax, psfext_ystep,
                    mc_ymin, mc_ymax, mc_ystep,
                    show_sdd11=show_sdd11, show_sdd22=show_sdd22,
                    next_files=next_files_tuple, fext_files=fext_files_tuple,
                )
            st.download_button(
                "📥 下載 PPTX",
                pptx_bytes,
                file_name=f"{stem}_report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        except Exception as e:
            st.error(f"需要 kaleido：pip install kaleido\n{e}")

        csv_bytes = _build_csv_zip(files_tuple, mapping, rise_time_ps)
        st.download_button(
            "📥 下載全部數據 (CSV)",
            csv_bytes,
            file_name=f"{stem}_data.zip",
            mime="application/zip",
        )

    next_fd_datasets = []
    fext_fd_datasets = []
    for uf in (uploaded_next_files or []):
        fb = uf.getvalue()
        with tempfile.NamedTemporaryFile(delete=False, suffix=".s4p") as tmp:
            tmp.write(fb); tmp_paths.append(tmp.name)
        fd, _ = _process_file(fb, mapping, rise_time_ps, uf.name, tmp_paths[-1])
        next_fd_datasets.append(fd)
    for uf in (uploaded_fext_files or []):
        fb = uf.getvalue()
        with tempfile.NamedTemporaryFile(delete=False, suffix=".s4p") as tmp:
            tmp.write(fb); tmp_paths.append(tmp.name)
        fd, _ = _process_file(fb, mapping, rise_time_ps, uf.name, tmp_paths[-1])
        fext_fd_datasets.append(fd)

    col1, col2 = st.columns(2)
    with col1:
        st.plotly_chart(plot_insertion_loss(
            fd_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
            y_min=il_ymin, y_max=il_ymax, y_step=il_ystep,
        ), use_container_width=True)
        st.plotly_chart(plot_psnext(
            next_fd_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
            y_min=psnext_ymin, y_max=psnext_ymax, y_step=psnext_ystep,
        ), use_container_width=True)
        st.plotly_chart(plot_mode_conversion(
            fd_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
            y_min=mc_ymin, y_max=mc_ymax, y_step=mc_ystep,
        ), use_container_width=True)

    with col2:
        st.plotly_chart(plot_return_loss(
            fd_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
            y_min=rl_ymin, y_max=rl_ymax, y_step=rl_ystep,
            show_sdd11=show_sdd11, show_sdd22=show_sdd22,
        ), use_container_width=True)
        st.plotly_chart(plot_psfext(
            fext_fd_datasets, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
            y_min=psfext_ymin, y_max=psfext_ymax, y_step=psfext_ystep,
        ), use_container_width=True)
        st.plotly_chart(plot_impedance(
            td_datasets, rise_time_ps=rise_time_ps,
            show_zdiff=show_zdiff, show_zse=show_zse,
            show_forward=show_forward, show_reverse=show_reverse,
            x_min=z_xmin, x_max=z_xmax, x_step=z_xstep,
            y_min=z_ymin, y_max=z_ymax, y_step=z_ystep,
        ), use_container_width=True)

except Exception as e:
    st.error(f"讀取失敗：{e}")

finally:
    for p in tmp_paths:
        try:
            os.unlink(p)
        except Exception:
            pass
