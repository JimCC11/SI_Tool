import io
import os
import tempfile
import zipfile

import numpy as np
import streamlit as st

from core.impedance import compute_tdr_diff, compute_tdr_single
from core.mixed_mode import get_pair_fd, single_to_mixed_mode_npairs
from core.parser import get_frequency_ghz, get_port_z0, load_snp
from core.plots import (
    plot_impedance,
    plot_insertion_loss,
    plot_mode_conversion,
    plot_psfext,
    plot_psnext,
    plot_return_loss,
)

_SNP_TYPES = ["s4p", "s8p", "s12p", "s16p", "s20p", "s24p", "s28p", "s32p"]
_IMG_EXPORT = dict(format="png", width=630, height=450, scale=2)
_TABLE_STYLE = """
<style>
.pm-table { width:100%; border-collapse:collapse; text-align:center; font-size:0.85rem; }
.pm-table td { padding: 4px 6px; border: 1px solid #444; }
.pm-diff { font-weight:bold; background:#dbeafe; color:#1e40af; }
</style>
"""


def _n_pairs_from_filename(filename: str) -> int:
    """從副檔名解析 port 數，例如 's16p' → 4 pairs。"""
    try:
        ext = filename.rsplit(".", 1)[-1].lower()   # e.g. "s16p"
        n_ports = int(ext[1:-1])                    # strip 's' and 'p'
        assert n_ports % 4 == 0 and 4 <= n_ports <= 32
        return n_ports // 4
    except Exception:
        return 1


def _port_map_html(n_pairs: int, mapping: str) -> str:
    rows = [_TABLE_STYLE + '<table class="pm-table">']
    for k in range(n_pairs):
        p1, p2, p3, p4 = 4*k+1, 4*k+2, 4*k+3, 4*k+4
        d1, d2 = 2*k+1, 2*k+2
        if mapping == 'A':
            a1, a2, b1, b2 = p1, p2, p3, p4
        else:
            a1, a2, b1, b2 = p1, p3, p2, p4
        rows.append(
            f'<tr><td class="pm-diff" rowspan="2">Diff {d1}</td>'
            f'<td>Port {a1}</td><td>────</td><td>Port {a2}</td>'
            f'<td class="pm-diff" rowspan="2">Diff {d2}</td></tr>'
            f'<tr><td>Port {b1}</td><td>────</td><td>Port {b2}</td></tr>'
        )
    rows.append("</table>")
    return "".join(rows)


def _process_snp(tmp_path: str, n_pairs: int, mapping: str, rise_time_ps: float):
    net     = load_snp(tmp_path)
    freq    = get_frequency_ghz(net)
    fhz     = net.f
    s_se    = net.s
    s_mm    = single_to_mixed_mode_npairs(s_se, n_pairs, mapping=mapping)
    z0_se   = get_port_z0(net)
    z0_diff = 2.0 * z0_se

    pair_labels = [f"Pair {i+1}" for i in range(n_pairs)]
    fd_list, td_list = [], []
    for i in range(n_pairs):
        fd = get_pair_fd(s_mm, freq, i, pair_labels[i])
        idx   = [4*i, 4*i+1, 4*i+2, 4*i+3]
        s_sub = s_se[:, idx, :][:, :, idx]
        t_f,  z11_f  = compute_tdr_single(s_sub, fhz, rise_time_ps=rise_time_ps, forward=True,  z0_se=z0_se)
        _,    zdif_f = compute_tdr_diff(s_mm[:, 4*i:4*i+4, 4*i:4*i+4], fhz, rise_time_ps=rise_time_ps, forward=True,  z0_diff=z0_diff)
        t_r,  z11_r  = compute_tdr_single(s_sub, fhz, rise_time_ps=rise_time_ps, forward=False, z0_se=z0_se)
        _,    zdif_r = compute_tdr_diff(s_mm[:, 4*i:4*i+4, 4*i:4*i+4], fhz, rise_time_ps=rise_time_ps, forward=False, z0_diff=z0_diff)
        td_list.append({
            "t_fwd": t_f,  "z11_fwd":  z11_f,  "zdiff_fwd": zdif_f,
            "t_rev": t_r,  "z11_rev":  z11_r,  "zdiff_rev": zdif_r,
            "label": pair_labels[i],
        })
        fd_list.append(fd)

    return freq, fhz, s_mm, fd_list, td_list, z0_se




# ── Page config ────────────────────────────────────────────────────────────────
st.set_page_config(page_title="SI Tool — SNP Analyzer", page_icon="📡", layout="wide")
st.title("SI Tool — SNP Analyzer (S4P～S32P)")

# ── Sidebar ────────────────────────────────────────────────────────────────────
with st.sidebar:
    st.header("上傳檔案")
    uploaded = st.file_uploader(
        "選擇 SNP 檔案（S4P～S32P）",
        type=_SNP_TYPES,
        accept_multiple_files=True,
        key="snp_up",
    )

    n_pairs = _n_pairs_from_filename(uploaded[0].name) if uploaded else 1
    n_ports = n_pairs * 4
    tx_labels = [f"Diff {2*k+1}" for k in range(n_pairs)]

    if uploaded:
        st.caption(f"偵測：{n_ports}-port → {n_pairs} 差分對")

    st.divider()
    st.subheader("Port Mapping")
    mapping_choice = st.radio("Port Mapping", ["Odd-Even", "N+1"], index=0,
                              horizontal=True, label_visibility="hidden")
    mapping = 'A' if mapping_choice == 'Odd-Even' else 'B'
    st.markdown(_port_map_html(n_pairs, mapping), unsafe_allow_html=True)

    st.divider()
    show_legend = st.checkbox("顯示圖例", value=True, key="snp_legend")

    st.divider()
    st.subheader("Frequency Domain")
    cx1, cx2, cx3 = st.columns(3)
    with cx1: fd_xmin  = st.number_input("X min (GHz)",  value=0.0,  step=1.0,            key="snp_fxn")
    with cx2: fd_xmax  = st.number_input("X max (GHz)",  value=20.0, step=1.0,            key="snp_fxx")
    with cx3: fd_xstep = st.number_input("X step (GHz)", value=5.0,  step=1.0, min_value=0.1, key="snp_fxs")

    st.caption("Insertion Loss")
    iy1, iy2, iy3 = st.columns(3)
    with iy1: il_ymin  = st.number_input("Y min",  key="snp_iln", value=-5.0,  step=1.0)
    with iy2: il_ymax  = st.number_input("Y max",  key="snp_ilx", value=0.0,   step=1.0)
    with iy3: il_ystep = st.number_input("Y step", key="snp_ils", value=1.0,   step=0.5, min_value=0.1)

    st.caption("Return Loss")
    ry1, ry2, ry3 = st.columns(3)
    with ry1: rl_ymin  = st.number_input("Y min",  key="snp_rln", value=-80.0, step=10.0)
    with ry2: rl_ymax  = st.number_input("Y max",  key="snp_rlx", value=0.0,   step=10.0)
    with ry3: rl_ystep = st.number_input("Y step", key="snp_rls", value=20.0,  step=5.0, min_value=0.1)
    rl_col1, rl_col2 = st.columns(2)
    with rl_col1: show_sdd11 = st.checkbox("SDD11", value=True,  key="snp_s11")
    with rl_col2: show_sdd22 = st.checkbox("SDD22", value=True,  key="snp_s22")

    st.caption("PSNEXT")
    nx1, nx2, nx3 = st.columns(3)
    with nx1: psnext_ymin  = st.number_input("Y min",  key="snp_nxn", value=-80.0, step=10.0)
    with nx2: psnext_ymax  = st.number_input("Y max",  key="snp_nxx", value=0.0,   step=10.0)
    with nx3: psnext_ystep = st.number_input("Y step", key="snp_nxs", value=20.0,  step=5.0, min_value=0.1)
    _mf = len(uploaded) > 1 if uploaded else False
    if uploaded and n_pairs >= 2:
        _nx_opts = [
            (f"{fi+1}, Diff{2*v+1}, Diff{2*a+1}" if _mf else f"Diff{2*v+1}, Diff{2*a+1}")
            for fi in range(len(uploaded))
            for v in range(n_pairs)
            for a in range(n_pairs) if a != v
        ]
        psnext_path_sel = st.multiselect("路徑", _nx_opts, default=_nx_opts, key="snp_nx_paths")
    else:
        psnext_path_sel = []

    st.caption("PSFEXT")
    fx1, fx2, fx3 = st.columns(3)
    with fx1: psfext_ymin  = st.number_input("Y min",  key="snp_fxn2", value=-80.0, step=10.0)
    with fx2: psfext_ymax  = st.number_input("Y max",  key="snp_fxx2", value=0.0,   step=10.0)
    with fx3: psfext_ystep = st.number_input("Y step", key="snp_fxs2", value=20.0,  step=5.0, min_value=0.1)
    if uploaded and n_pairs >= 2:
        _fx_opts = [
            (f"{fi+1}, Diff{2*v+1}, Diff{2*a+2}" if _mf else f"Diff{2*v+1}, Diff{2*a+2}")
            for fi in range(len(uploaded))
            for v in range(n_pairs)
            for a in range(n_pairs) if a != v
        ]
        psfext_path_sel = st.multiselect("路徑", _fx_opts, default=_fx_opts, key="snp_fx_paths")
    else:
        psfext_path_sel = []

    st.caption("Mode Conversion")
    my1, my2, my3 = st.columns(3)
    with my1: mc_ymin  = st.number_input("Y min",  key="snp_mcn", value=-80.0, step=10.0)
    with my2: mc_ymax  = st.number_input("Y max",  key="snp_mcx", value=0.0,   step=10.0)
    with my3: mc_ystep = st.number_input("Y step", key="snp_mcs", value=20.0,  step=5.0, min_value=0.1)

    st.divider()
    st.subheader("Time Domain")
    rise_time_ps = st.number_input("Rise time 20%~80% (ps)", min_value=10, max_value=500,
                                   value=35, step=1, key="snp_rt")
    st.caption(f"建議頻寬：{0.35 / (rise_time_ps * 1e-12) / 1e9 * 2:.1f} GHz")
    col_xmin, col_xmax, col_xstep = st.columns(3)
    with col_xmin:  z_xmin  = st.number_input("X min (ns)",  value=0.0, step=0.1, format="%.2f",  key="snp_zxn")
    with col_xmax:  z_xmax  = st.number_input("X max (ns)",  value=1.0, step=0.1, format="%.2f",  key="snp_zxx")
    with col_xstep: z_xstep = st.number_input("X step (ns)", value=0.2, step=0.1, min_value=0.001,
                                               format="%.3f", key="snp_zxs")
    col_l, col_r = st.columns(2)
    with col_l: show_zdiff = st.checkbox("Z_Diff", value=True,  key="snp_zd")
    with col_r: show_zse   = st.checkbox("Z_SE",   value=False, key="snp_zse")
    col_fwd, col_rev = st.columns(2)
    with col_fwd: show_forward = st.checkbox("Forward", value=True,  key="snp_fwd")
    with col_rev: show_reverse = st.checkbox("Reverse", value=False, key="snp_rev")
    col_min, col_max, col_ystep = st.columns(3)
    with col_min:   z_ymin  = st.number_input("Y min (Ω)",  value=0,   step=10, key="snp_zyn")
    with col_max:   z_ymax  = st.number_input("Y max (Ω)",  value=150, step=10, key="snp_zyx")
    with col_ystep: z_ystep = st.number_input("Y step (Ω)", value=50,  step=10, min_value=1, key="snp_zys")

# ── Main ───────────────────────────────────────────────────────────────────────
if not uploaded:
    st.info("← 請從左側上傳 SNP 檔案（S4P～S32P）")
    st.stop()

def _diff_tx(k): return 2 * k + 1
def _diff_rx(k): return 2 * k + 2

tmp_paths = []
try:
    multi_file = len(uploaded) > 1
    fd_list, td_list, _all_psnext, _all_psfext = [], [], [], []
    per_file_data = []
    info_lines    = []

    for fi, uf in enumerate(uploaded):
        stem   = uf.name.rsplit(".", 1)[0]
        suffix = "." + uf.name.rsplit(".", 1)[-1]
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(uf.getvalue())
            tmp_paths.append(tmp.name)

        freq, fhz, s_mm, fds, tds, z0_se = _process_snp(tmp_paths[-1], n_pairs, mapping, rise_time_ps)
        info_lines.append(f"**{uf.name}** {freq[0]:.2f}~{freq[-1]:.2f} GHz | Z0={z0_se:.1f}Ω (SE) / {2*z0_se:.1f}Ω (Diff)")

        pfx = f"{stem} " if multi_file else ""
        for k, fd in enumerate(fds):
            tx, rx = _diff_tx(k), _diff_rx(k)
            fd["trace_labels"] = {
                "sdd21": f"{pfx}SDD{rx}{tx}",
                "sdd11": f"{pfx}SDD{tx}{tx}",
                "sdd22": f"{pfx}SDD{rx}{rx}",
                "scd21": f"{pfx}SCD{rx}{tx}",
                "sdc21": f"{pfx}SDC{rx}{tx}",
            }
        per_file_data.append((stem, freq, fds, tds))
        fd_list.extend(fds)
        td_list.extend(tds)

        if n_pairs >= 2:
            _all_psnext.extend([
                {"freq": freq, "sdd21": s_mm[:, 4*v, 4*a],
                 "label": (f"{fi+1}, Diff{_diff_tx(v)}, Diff{_diff_tx(a)}" if multi_file
                           else f"Diff{_diff_tx(v)}, Diff{_diff_tx(a)}"),
                 "_victim": v, "_agg": a, "_stem": stem, "_file_idx": fi}
                for v in range(n_pairs) for a in range(n_pairs) if a != v
            ])
            _all_psfext.extend([
                {"freq": freq, "sdd21": s_mm[:, 4*v+1, 4*a],
                 "label": (f"{fi+1}, Diff{_diff_tx(v)}, Diff{_diff_rx(a)}" if multi_file
                           else f"Diff{_diff_tx(v)}, Diff{_diff_rx(a)}"),
                 "_victim": v, "_agg": a, "_stem": stem, "_file_idx": fi}
                for v in range(n_pairs) for a in range(n_pairs) if a != v
            ])

    st.success(" | ".join(info_lines) + f" | {n_pairs} 對差分對")

    psnext_fd = [ds for ds in _all_psnext if ds["label"] in psnext_path_sel]
    psfext_fd = [ds for ds in _all_psfext if ds["label"] in psfext_path_sel]

    def _per_victim_ps(selected_ds, label_fn):
        groups = {}
        for ds in selected_ds:
            groups.setdefault((ds["_file_idx"], ds["_victim"]), []).append(ds)
        result = []
        for (fi, v), terms in sorted(groups.items()):
            ps  = np.sum([np.abs(ds["sdd21"])**2 for ds in terms], axis=0)
            pfx = f"{fi+1}, " if multi_file else ""
            result.append({"freq": terms[0]["freq"], "sdd21": np.sqrt(ps),
                           "label": f"{pfx}{label_fn(v)}"})
        return result

    psnext_ps = _per_victim_ps(psnext_fd, lambda v: f"PSNEXT Diff{_diff_tx(v)}")
    psfext_ps = _per_victim_ps(psfext_fd, lambda v: f"PSFEXT Diff{_diff_tx(v)}")

    def _fig(f):
        f.update_layout(showlegend=show_legend)
        return f

    export_stem = uploaded[0].name.rsplit(".", 1)[0] if len(uploaded) == 1 else "SNP"
    figs_export = {
        "01_insertion_loss":  _fig(plot_insertion_loss(
            fd_list, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
            y_min=il_ymin, y_max=il_ymax, y_step=il_ystep)),
        "02_return_loss":     _fig(plot_return_loss(
            fd_list, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
            y_min=rl_ymin, y_max=rl_ymax, y_step=rl_ystep,
            show_sdd11=show_sdd11, show_sdd22=show_sdd22)),
        "03_psnext":          _fig(plot_psnext(
            psnext_fd, ps_datasets=psnext_ps,
            x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
            y_min=psnext_ymin, y_max=psnext_ymax, y_step=psnext_ystep)),
        "04_psfext":          _fig(plot_psfext(
            psfext_fd, ps_datasets=psfext_ps,
            x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
            y_min=psfext_ymin, y_max=psfext_ymax, y_step=psfext_ystep)),
        "05_mode_conversion": _fig(plot_mode_conversion(
            fd_list, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep,
            y_min=mc_ymin, y_max=mc_ymax, y_step=mc_ystep)),
        "06_tdr_impedance":   _fig(plot_impedance(
            td_list, rise_time_ps=rise_time_ps,
            show_zdiff=show_zdiff, show_zse=show_zse,
            show_forward=show_forward, show_reverse=show_reverse,
            x_min=z_xmin, x_max=z_xmax, x_step=z_xstep,
            y_min=z_ymin, y_max=z_ymax, y_step=z_ystep)),
    }

    # ── 下載 ──────────────────────────────────────────────────────────────────
    with st.sidebar:
        st.divider()
        st.subheader("輸出")

        try:
            with st.spinner("生成圖片中…"):
                buf = io.BytesIO()
                with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
                    for name, fig in figs_export.items():
                        zf.writestr(f"{name}.png", fig.to_image(**_IMG_EXPORT))
            st.download_button("📥 下載全部圖表 (PNG)", buf.getvalue(),
                               file_name=f"{export_stem}_charts.zip", mime="application/zip")
        except Exception as e:
            st.error(f"需要 kaleido：pip install kaleido\n{e}")

        try:
            from pptx import Presentation
            template_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "Template.pptx")
            prs = Presentation(template_path)
            slide_map = [
                {13: figs_export["01_insertion_loss"], 14: figs_export["02_return_loss"],
                 15: figs_export["03_psnext"],         16: figs_export["04_psfext"]},
                {13: figs_export["05_mode_conversion"]},
                {13: figs_export["06_tdr_impedance"]},
            ]
            for slide, fmap in zip(prs.slides, slide_map):
                for ph in slide.placeholders:
                    idx = ph.placeholder_format.idx
                    if idx in fmap:
                        ph.insert_picture(io.BytesIO(fmap[idx].to_image(**_IMG_EXPORT)))
            pptx_buf = io.BytesIO()
            prs.save(pptx_buf)
            st.download_button("📥 下載 PPTX", pptx_buf.getvalue(),
                               file_name=f"{export_stem}_report.pptx",
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        except ImportError:
            st.warning("需要 python-pptx：pip install python-pptx")
        except Exception as e:
            st.error(f"PPTX 失敗：{e}")

        try:
            csv_buf = io.BytesIO()
            with zipfile.ZipFile(csv_buf, "w", zipfile.ZIP_DEFLATED) as zf:
                for stem_c, freq_c, fds_c, tds_c in per_file_data:
                    hdr = ["Frequency_GHz"]
                    for k in range(n_pairs):
                        tx, rx = _diff_tx(k), _diff_rx(k)
                        hdr += [f"SDD{rx}{tx}_dB", f"SDD{tx}{tx}_dB", f"SDD{rx}{rx}_dB",
                                f"SCD{rx}{tx}_dB", f"SDC{rx}{tx}_dB"]
                    fd_rows = [",".join(hdr)]
                    for i, fv in enumerate(freq_c):
                        vals = [f"{fv:.6f}"]
                        for fd in fds_c:
                            for key in ["sdd21", "sdd11", "sdd22", "scd21", "sdc21"]:
                                vals.append(f"{20*np.log10(np.abs(fd[key][i])+1e-15):.6f}")
                        fd_rows.append(",".join(vals))
                    zf.writestr(f"{stem_c}_FD.csv", "\n".join(fd_rows))

                    hdr_td = ["Time_ns"]
                    for k in range(n_pairs):
                        tx = _diff_tx(k)
                        hdr_td += [f"ZSE_Fwd_Diff{tx}_Ohm", f"ZDiff_Fwd_Diff{tx}_Ohm",
                                   f"ZSE_Rev_Diff{tx}_Ohm",  f"ZDiff_Rev_Diff{tx}_Ohm"]
                    td_rows = [",".join(hdr_td)]
                    n_pts = len(tds_c[0]["t_fwd"])
                    for i in range(n_pts):
                        vals = [f"{tds_c[0]['t_fwd'][i]/1000:.6f}"]
                        for td in tds_c:
                            vals += [f"{td['z11_fwd'][i]:.6f}", f"{td['zdiff_fwd'][i]:.6f}",
                                     f"{td['z11_rev'][i]:.6f}",  f"{td['zdiff_rev'][i]:.6f}"]
                        td_rows.append(",".join(vals))
                    zf.writestr(f"{stem_c}_TD.csv", "\n".join(td_rows))
            st.download_button("📥 下載全部數據 (CSV)", csv_buf.getvalue(),
                               file_name=f"{export_stem}_data.zip", mime="application/zip")
        except Exception as e:
            st.error(f"CSV 失敗：{e}")

    # ── 版面 3×2 ──────────────────────────────────────────────────────────────
    col1, col2 = st.columns(2)
    with col1:
        st.plotly_chart(figs_export["01_insertion_loss"],  use_container_width=True)
        st.plotly_chart(figs_export["03_psnext"],          use_container_width=True)
        st.plotly_chart(figs_export["05_mode_conversion"], use_container_width=True)
    with col2:
        st.plotly_chart(figs_export["02_return_loss"],     use_container_width=True)
        st.plotly_chart(figs_export["04_psfext"],          use_container_width=True)
        st.plotly_chart(figs_export["06_tdr_impedance"],   use_container_width=True)

except Exception as e:
    st.error(f"讀取失敗：{e}")
finally:
    for p in tmp_paths:
        try:
            os.unlink(p)
        except Exception:
            pass
