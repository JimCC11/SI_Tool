import io
import os
import tempfile
import zipfile

import numpy as np
import streamlit as st

from core.impedance import compute_tdr_diff, compute_tdr_single
from core.mixed_mode import (
    get_pair_fd,
    single_to_mixed_mode_npairs,
)
from core.parser import get_frequency_ghz, load_s16p
from core.plots import (
    plot_impedance,
    plot_insertion_loss,
    plot_mode_conversion,
    plot_psfext,
    plot_psnext,
    plot_return_loss,
)

N_PAIRS = 4
PAIR_LABELS = [f"Pair {i+1}" for i in range(N_PAIRS)]
_IMG_EXPORT = dict(format="png", width=630, height=450, scale=2)


def _process_s16p(tmp_path: str, mapping: str, rise_time_ps: float):
    net    = load_s16p(tmp_path)
    freq   = get_frequency_ghz(net)
    fhz    = net.f
    s_se   = net.s
    s_mm   = single_to_mixed_mode_npairs(s_se, N_PAIRS, mapping=mapping)

    fd_list, td_list = [], []
    for i in range(N_PAIRS):
        fd = get_pair_fd(s_mm, freq, i, PAIR_LABELS[i])
        # extract 4x4 sub-matrix for this pair
        idx = [4*i, 4*i+1, 4*i+2, 4*i+3]
        s_sub = s_se[:, np.ix_(idx, idx)[0], np.ix_(idx, idx)[1]] if False else s_se[:, idx, :][:, :, idx]
        t_f, z11_f = compute_tdr_single(s_sub, fhz, rise_time_ps=rise_time_ps, forward=True)
        _,   zdif_f = compute_tdr_diff(s_mm[:, 4*i:4*i+4, 4*i:4*i+4], fhz, rise_time_ps=rise_time_ps, forward=True)
        t_r, z11_r = compute_tdr_single(s_sub, fhz, rise_time_ps=rise_time_ps, forward=False)
        _,   zdif_r = compute_tdr_diff(s_mm[:, 4*i:4*i+4, 4*i:4*i+4], fhz, rise_time_ps=rise_time_ps, forward=False)
        td = {
            "t_fwd": t_f, "z11_fwd": z11_f, "zdiff_fwd": zdif_f,
            "t_rev": t_r, "z11_rev": z11_r, "zdiff_rev": zdif_r,
            "label": PAIR_LABELS[i],
        }
        fd_list.append(fd)
        td_list.append(td)

    return freq, fhz, s_mm, fd_list, td_list


# ── Page config ───────────────────────────────────────────
st.set_page_config(page_title="SI Tool — S16P", page_icon="📡", layout="wide")
st.title("SI Tool — S16P Analyzer")

# ── Sidebar ──────────────────────────────────────────────
with st.sidebar:
    st.header("上傳檔案")
    uploaded = st.file_uploader("選擇 S16P 檔案", type=["s16p", "s4p"], accept_multiple_files=True)

    st.divider()
    st.subheader("Port Mapping")
    mapping_choice = st.radio("Port Mapping", ["Odd-Even", "N+1"], index=0,
                              horizontal=True, label_visibility="hidden")
    mapping = 'A' if mapping_choice == 'Odd-Even' else 'B'

    _TABLE_STYLE = """
<style>
.pm-table { width:100%; border-collapse:collapse; text-align:center; font-size:0.85rem; }
.pm-table td { padding: 4px 6px; border: 1px solid #444; }
.pm-diff { font-weight:bold; background:#dbeafe; color:#1e40af; }
</style>
"""
    if mapping == 'A':
        st.markdown(_TABLE_STYLE + """
<table class="pm-table">
  <tr>
    <td class="pm-diff" rowspan="2">Diff 1</td>
    <td>Port 1</td><td>────</td><td>Port 2</td>
    <td class="pm-diff" rowspan="2">Diff 2</td>
  </tr>
  <tr><td>Port 3</td><td>────</td><td>Port 4</td></tr>
  <tr>
    <td class="pm-diff" rowspan="2">Diff 3</td>
    <td>Port 5</td><td>────</td><td>Port 6</td>
    <td class="pm-diff" rowspan="2">Diff 4</td>
  </tr>
  <tr><td>Port 7</td><td>────</td><td>Port 8</td></tr>
  <tr>
    <td class="pm-diff" rowspan="2">Diff 5</td>
    <td>Port 9</td><td>────</td><td>Port 10</td>
    <td class="pm-diff" rowspan="2">Diff 6</td>
  </tr>
  <tr><td>Port 11</td><td>────</td><td>Port 12</td></tr>
  <tr>
    <td class="pm-diff" rowspan="2">Diff 7</td>
    <td>Port 13</td><td>────</td><td>Port 14</td>
    <td class="pm-diff" rowspan="2">Diff 8</td>
  </tr>
  <tr><td>Port 15</td><td>────</td><td>Port 16</td></tr>
</table>""", unsafe_allow_html=True)
    else:
        st.markdown(_TABLE_STYLE + """
<table class="pm-table">
  <tr>
    <td class="pm-diff" rowspan="2">Diff 1</td>
    <td>Port 1</td><td>────</td><td>Port 3</td>
    <td class="pm-diff" rowspan="2">Diff 2</td>
  </tr>
  <tr><td>Port 2</td><td>────</td><td>Port 4</td></tr>
  <tr>
    <td class="pm-diff" rowspan="2">Diff 3</td>
    <td>Port 5</td><td>────</td><td>Port 7</td>
    <td class="pm-diff" rowspan="2">Diff 4</td>
  </tr>
  <tr><td>Port 6</td><td>────</td><td>Port 8</td></tr>
  <tr>
    <td class="pm-diff" rowspan="2">Diff 5</td>
    <td>Port 9</td><td>────</td><td>Port 11</td>
    <td class="pm-diff" rowspan="2">Diff 6</td>
  </tr>
  <tr><td>Port 10</td><td>────</td><td>Port 12</td></tr>
  <tr>
    <td class="pm-diff" rowspan="2">Diff 7</td>
    <td>Port 13</td><td>────</td><td>Port 15</td>
    <td class="pm-diff" rowspan="2">Diff 8</td>
  </tr>
  <tr><td>Port 14</td><td>────</td><td>Port 16</td></tr>
</table>""", unsafe_allow_html=True)

    st.divider()
    show_legend = st.checkbox("顯示圖例", value=True, key="show_legend")
    st.divider()
    st.subheader("Frequency Domain")
    cx1, cx2, cx3 = st.columns(3)
    with cx1: fd_xmin  = st.number_input("X min (GHz)", value=0.0,  step=1.0)
    with cx2: fd_xmax  = st.number_input("X max (GHz)", value=20.0, step=1.0)
    with cx3: fd_xstep = st.number_input("X step (GHz)", value=5.0, step=1.0, min_value=0.1)

    st.caption("Insertion Loss")
    iy1, iy2, iy3 = st.columns(3)
    with iy1: il_ymin  = st.number_input("Y min", key="il_ymin",  value=-5.0,  step=1.0)
    with iy2: il_ymax  = st.number_input("Y max", key="il_ymax",  value=0.0,   step=1.0)
    with iy3: il_ystep = st.number_input("Y step", key="il_ystep", value=1.0,  step=0.5, min_value=0.1)

    st.caption("Return Loss")
    ry1, ry2, ry3 = st.columns(3)
    with ry1: rl_ymin  = st.number_input("Y min", key="rl_ymin",  value=-80.0, step=10.0)
    with ry2: rl_ymax  = st.number_input("Y max", key="rl_ymax",  value=0.0,   step=10.0)
    with ry3: rl_ystep = st.number_input("Y step", key="rl_ystep", value=20.0, step=5.0, min_value=0.1)
    rl_col1, rl_col2 = st.columns(2)
    with rl_col1: show_sdd11 = st.checkbox("SDD11", value=True)
    with rl_col2: show_sdd22 = st.checkbox("SDD22", value=True)

    _TX_LABELS = [f"Diff {2*k+1}" for k in range(N_PAIRS)]  # ["Diff 1","Diff 3","Diff 5","Diff 7"]

    def _victim_agg_ui(prefix: str, y_keys: tuple):
        """動態 Victim 列表：每個 Victim 各自有 Aggressors multiselect。
        最後一列旁邊有 ＋ 新增 / － 移除 按鍵。
        回傳 (ymin, ymax, ystep, {victim_pair_idx: set_of_agg_pair_idx})。
        """
        yk_min, yk_max, yk_step = y_keys
        c1, c2, c3 = st.columns(3)
        with c1: ymin  = st.number_input("Y min",  key=yk_min,  value=-80.0, step=10.0)
        with c2: ymax  = st.number_input("Y max",  key=yk_max,  value=0.0,   step=10.0)
        with c3: ystep = st.number_input("Y step", key=yk_step, value=20.0,  step=5.0, min_value=0.1)

        n_key = f"{prefix}_n_victims"
        if n_key not in st.session_state:
            st.session_state[n_key] = 1
        n = st.session_state[n_key]

        sel = {}
        for i in range(n):
            v_key = f"{prefix}_victim_{i}"
            a_key = f"{prefix}_aggs_{i}"
            is_last = (i == n - 1)

            if is_last:
                col_v, col_add, col_rm = st.columns([3, 1, 1])
            else:
                col_v = st.columns(1)[0]

            with col_v:
                cur_v = st.selectbox(f"Victim_{i+1}", _TX_LABELS, key=v_key)
            if is_last:
                with col_add:
                    st.write("")
                    if n < N_PAIRS and st.button("＋", key=f"{prefix}_add"):
                        st.session_state[n_key] += 1
                        st.rerun()
                with col_rm:
                    st.write("")
                    if n > 1 and st.button("－", key=f"{prefix}_rm"):
                        st.session_state[n_key] -= 1
                        st.rerun()

            cur_k = _TX_LABELS.index(cur_v)
            agg_opts = [_TX_LABELS[m] for m in range(N_PAIRS) if m != cur_k]
            chosen = st.multiselect("Aggressors", agg_opts, default=agg_opts, key=a_key)
            sel[cur_k] = {_TX_LABELS.index(l) for l in chosen}

        return ymin, ymax, ystep, sel

    st.caption("PSNEXT")
    psnext_ymin, psnext_ymax, psnext_ystep, psnext_sel = _victim_agg_ui(
        "psnext", ("psnext_ymin", "psnext_ymax", "psnext_ystep"))

    st.caption("PSFEXT")
    psfext_ymin, psfext_ymax, psfext_ystep, psfext_sel = _victim_agg_ui(
        "psfext", ("psfext_ymin", "psfext_ymax", "psfext_ystep"))

    st.caption("Mode Conversion")
    my1, my2, my3 = st.columns(3)
    with my1: mc_ymin  = st.number_input("Y min", key="mc_ymin",  value=-80.0, step=10.0)
    with my2: mc_ymax  = st.number_input("Y max", key="mc_ymax",  value=0.0,   step=10.0)
    with my3: mc_ystep = st.number_input("Y step", key="mc_ystep", value=20.0, step=5.0, min_value=0.1)

    st.divider()
    st.subheader("Time Domain")
    rise_time_ps = st.number_input("Rise time 20%~80% (ps)", min_value=10, max_value=500, value=35, step=1)
    st.caption(f"建議頻寬：{0.35 / (rise_time_ps * 1e-12) / 1e9 * 2:.1f} GHz")
    col_xmin, col_xmax, col_xstep = st.columns(3)
    with col_xmin:  z_xmin  = st.number_input("X min (ns)", value=0.0, step=0.1, format="%.2f")
    with col_xmax:  z_xmax  = st.number_input("X max (ns)", value=1.0, step=0.1, format="%.2f")
    with col_xstep: z_xstep = st.number_input("X step (ns)", value=0.2, step=0.1, min_value=0.001, format="%.3f")
    col_l, col_r = st.columns(2)
    with col_l: show_zdiff = st.checkbox("Z_Diff", value=True)
    with col_r: show_zse   = st.checkbox("Z_SE",   value=False)
    col_fwd, col_rev = st.columns(2)
    with col_fwd: show_forward = st.checkbox("Forward", value=True)
    with col_rev: show_reverse = st.checkbox("Reverse", value=False)
    col_min, col_max, col_ystep = st.columns(3)
    with col_min:   z_ymin  = st.number_input("Y min (Ω)", value=0,   step=10)
    with col_max:   z_ymax  = st.number_input("Y max (Ω)", value=150, step=10)
    with col_ystep: z_ystep = st.number_input("Y step (Ω)", value=50, step=10, min_value=1)

# ── Main ─────────────────────────────────────────────────
if not uploaded:
    st.info("← 請從左側上傳 S16P 檔案")
    st.stop()

def _diff_tx(k): return 2 * k + 1
def _diff_rx(k): return 2 * k + 2

tmp_paths = []
try:
    multi_file = len(uploaded) > 1
    fd_list, td_list, _all_psnext, _all_psfext = [], [], [], []
    per_file_data = []   # [(stem, freq, fds, tds), ...]
    info_lines = []

    for uf in uploaded:
        stem = uf.name.rsplit(".", 1)[0]
        with tempfile.NamedTemporaryFile(delete=False, suffix=".s16p") as tmp:
            tmp.write(uf.getvalue())
            tmp_paths.append(tmp.name)

        freq, fhz, s_mm, fds, tds = _process_s16p(tmp_paths[-1], mapping, rise_time_ps)
        info_lines.append(f"**{uf.name}** {freq[0]:.2f}~{freq[-1]:.2f} GHz")

        pfx = f"{stem} " if multi_file else ""
        for k, fd in enumerate(fds):
            tx, rx = _diff_tx(k), _diff_rx(k)
            fd["trace_labels"] = {
                "sdd21": f"{pfx}SDD{rx}{tx}",
                "sdd11": f"{pfx}SDD{tx}{tx}",
                "sdd22": f"{pfx}SDD{rx}{rx}",
                "scd21": f"{pfx}SCD{rx}{tx}",
                "sdc21": f"{pfx}SDC{rx}{tx}",
            }
        per_file_data.append((stem, freq, fds, tds))
        fd_list.extend(fds)
        td_list.extend(tds)

        _all_psnext.extend([
            {"freq": freq, "sdd21": s_mm[:, 4*v, 4*a],
             "label": f"{pfx}SDD{_diff_tx(a)}{_diff_tx(v)}", "_victim": v, "_agg": a, "_stem": stem}
            for v in range(N_PAIRS) for a in range(N_PAIRS) if a != v
        ])
        _all_psfext.extend([
            {"freq": freq, "sdd21": s_mm[:, 4*v+1, 4*a],
             "label": f"{pfx}SDD{_diff_tx(a)}{_diff_rx(v)}", "_victim": v, "_agg": a, "_stem": stem}
            for v in range(N_PAIRS) for a in range(N_PAIRS) if a != v
        ])

    st.success(" | ".join(info_lines) + f" | {N_PAIRS} 對差分對")

    # 依每個 Victim 各自的 Aggressors 篩選個別虛線
    psnext_fd = [ds for ds in _all_psnext
                 if ds["_agg"] in psnext_sel.get(ds["_victim"], set())]
    psfext_fd = [ds for ds in _all_psfext
                 if ds["_agg"] in psfext_sel.get(ds["_victim"], set())]

    # 每個 (file, victim) 的 PSNEXT/PSFEXT 實線
    def _per_victim_ps(all_ds, sel, label_fn):
        stems = list(dict.fromkeys(ds["_stem"] for ds in all_ds))
        result = []
        for stem in stems:
            pfx = f"{stem} " if multi_file else ""
            for v, aggs in sorted(sel.items()):
                if not aggs:
                    continue
                terms = [ds for ds in all_ds
                         if ds["_victim"] == v and ds["_agg"] in aggs and ds["_stem"] == stem]
                if terms:
                    ps = np.sum([np.abs(ds["sdd21"])**2 for ds in terms], axis=0)
                    result.append({"freq": terms[0]["freq"], "sdd21": np.sqrt(ps),
                                   "label": f"{pfx}{label_fn(v)}"})
        return result

    psnext_ps = _per_victim_ps(_all_psnext, psnext_sel, lambda v: f"PSNEXT Diff{_diff_tx(v)}")
    psfext_ps = _per_victim_ps(_all_psfext, psfext_sel, lambda v: f"PSFEXT Diff{_diff_rx(v)}")

    def _fig(f):
        """套用圖例顯示設定"""
        f.update_layout(showlegend=show_legend)
        return f

    # ── 共用圖表物件 ──────────────────────────────────────
    export_stem = uploaded[0].name.rsplit(".", 1)[0] if len(uploaded) == 1 else "S16P"
    figs_export = {
        "01_insertion_loss":  _fig(plot_insertion_loss(fd_list, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep, y_min=il_ymin, y_max=il_ymax, y_step=il_ystep)),
        "02_return_loss":     _fig(plot_return_loss(fd_list, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep, y_min=rl_ymin, y_max=rl_ymax, y_step=rl_ystep, show_sdd11=show_sdd11, show_sdd22=show_sdd22)),
        "03_psnext":          _fig(plot_psnext(psnext_fd, ps_datasets=psnext_ps, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep, y_min=psnext_ymin, y_max=psnext_ymax, y_step=psnext_ystep)),
        "04_psfext":          _fig(plot_psfext(psfext_fd, ps_datasets=psfext_ps, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep, y_min=psfext_ymin, y_max=psfext_ymax, y_step=psfext_ystep)),
        "05_mode_conversion": _fig(plot_mode_conversion(fd_list, x_min=fd_xmin, x_max=fd_xmax, x_step=fd_xstep, y_min=mc_ymin, y_max=mc_ymax, y_step=mc_ystep)),
        "06_tdr_impedance":   _fig(plot_impedance(td_list, rise_time_ps=rise_time_ps, show_zdiff=show_zdiff, show_zse=show_zse, show_forward=show_forward, show_reverse=show_reverse, x_min=z_xmin, x_max=z_xmax, x_step=z_xstep, y_min=z_ymin, y_max=z_ymax, y_step=z_ystep)),
    }

    # ── 下載 ─────────────────────────────────────────────
    with st.sidebar:
        st.divider()
        st.subheader("輸出")

        # PNG ZIP
        try:
            with st.spinner("生成圖片中…"):
                buf = io.BytesIO()
                with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
                    for name, fig in figs_export.items():
                        zf.writestr(f"{name}.png", fig.to_image(**_IMG_EXPORT))
            st.download_button("📥 下載全部圖表 (PNG)", buf.getvalue(),
                               file_name=f"{export_stem}_charts.zip", mime="application/zip")
        except Exception as e:
            st.error(f"需要 kaleido：pip install kaleido\n{e}")

        # PPTX
        try:
            from pptx import Presentation
            template_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "Template.pptx")
            prs = Presentation(template_path)
            slide_map = [
                {13: figs_export["01_insertion_loss"], 14: figs_export["02_return_loss"],
                 15: figs_export["03_psnext"],         16: figs_export["04_psfext"]},
                {13: figs_export["05_mode_conversion"]},
                {13: figs_export["06_tdr_impedance"]},
            ]
            for slide, fmap in zip(prs.slides, slide_map):
                for ph in slide.placeholders:
                    idx = ph.placeholder_format.idx
                    if idx in fmap:
                        ph.insert_picture(io.BytesIO(fmap[idx].to_image(**_IMG_EXPORT)))
            pptx_buf = io.BytesIO()
            prs.save(pptx_buf)
            st.download_button("📥 下載 PPTX", pptx_buf.getvalue(),
                               file_name=f"{export_stem}_report.pptx",
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        except ImportError:
            st.warning("需要 python-pptx：pip install python-pptx")
        except Exception as e:
            st.error(f"PPTX 失敗：{e}")

        # CSV
        try:
            csv_buf = io.BytesIO()
            with zipfile.ZipFile(csv_buf, "w", zipfile.ZIP_DEFLATED) as zf:
                for stem_c, freq_c, fds_c, tds_c in per_file_data:
                    # Frequency Domain CSV
                    hdr = ["Frequency_GHz"]
                    for k in range(N_PAIRS):
                        tx, rx = _diff_tx(k), _diff_rx(k)
                        hdr += [f"SDD{rx}{tx}_dB", f"SDD{tx}{tx}_dB", f"SDD{rx}{rx}_dB",
                                f"SCD{rx}{tx}_dB", f"SDC{rx}{tx}_dB"]
                    fd_rows = [",".join(hdr)]
                    for i, fv in enumerate(freq_c):
                        vals = [f"{fv:.6f}"]
                        for fd in fds_c:
                            for key in ["sdd21", "sdd11", "sdd22", "scd21", "sdc21"]:
                                vals.append(f"{20*np.log10(np.abs(fd[key][i])+1e-15):.6f}")
                        fd_rows.append(",".join(vals))
                    zf.writestr(f"{stem_c}_FD.csv", "\n".join(fd_rows))

                    # Time Domain CSV
                    hdr_td = ["Time_ns"]
                    for k in range(N_PAIRS):
                        tx, rx = _diff_tx(k), _diff_rx(k)
                        hdr_td += [f"ZSE_Fwd_Diff{tx}_Ohm", f"ZDiff_Fwd_Diff{tx}_Ohm",
                                   f"ZSE_Rev_Diff{tx}_Ohm",  f"ZDiff_Rev_Diff{tx}_Ohm"]
                    td_rows = [",".join(hdr_td)]
                    n_pts = len(tds_c[0]["t_fwd"])
                    for i in range(n_pts):
                        vals = [f"{tds_c[0]['t_fwd'][i]/1000:.6f}"]
                        for td in tds_c:
                            vals += [f"{td['z11_fwd'][i]:.6f}", f"{td['zdiff_fwd'][i]:.6f}",
                                     f"{td['z11_rev'][i]:.6f}",  f"{td['zdiff_rev'][i]:.6f}"]
                        td_rows.append(",".join(vals))
                    zf.writestr(f"{stem_c}_TD.csv", "\n".join(td_rows))

            st.download_button("📥 下載全部數據 (CSV)", csv_buf.getvalue(),
                               file_name=f"{export_stem}_data.zip", mime="application/zip")
        except Exception as e:
            st.error(f"CSV 失敗：{e}")

    # ── 版面 3×2 ─────────────────────────────────────────
    col1, col2 = st.columns(2)
    with col1:
        st.plotly_chart(figs_export["01_insertion_loss"],  use_container_width=True)
        st.plotly_chart(figs_export["03_psnext"],          use_container_width=True)
        st.plotly_chart(figs_export["05_mode_conversion"], use_container_width=True)
    with col2:
        st.plotly_chart(figs_export["02_return_loss"],     use_container_width=True)
        st.plotly_chart(figs_export["04_psfext"],          use_container_width=True)
        st.plotly_chart(figs_export["06_tdr_impedance"],   use_container_width=True)

except Exception as e:
    st.error(f"讀取失敗：{e}")
finally:
    for p in tmp_paths:
        try:
            os.unlink(p)
        except Exception:
            pass
